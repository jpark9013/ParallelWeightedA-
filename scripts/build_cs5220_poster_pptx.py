#!/usr/bin/env python3
"""Build CS 5220 Parallel Weighted A* poster as a single-slide widescreen .pptx."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor


def add_card(slide, left, top, width, height, title, body, title_size=13, body_size=9):
    """Light bordered rounded rectangle (poster box) with text on top."""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0xFA, 0xFA, 0xFA)
    card.line.color.rgb = RGBColor(0xC0, 0xC0, 0xC0)
    card.line.width = Pt(1.1)
    pad = Inches(0.08)
    add_textbox(
        slide,
        left + pad,
        top + pad,
        width - 2 * pad,
        height - 2 * pad,
        title,
        body,
        title_size=title_size,
        body_size=body_size,
    )


def add_textbox(slide, left, top, width, height, title, body, title_size=13, body_size=9):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0.06)
    p = tf.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(title_size)
    p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    if body:
        p2 = tf.add_paragraph()
        p2.text = body
        p2.font.size = Pt(body_size)
        p2.font.color.rgb = RGBColor(0x22, 0x22, 0x22)
        p2.space_before = Pt(4)
    return box


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    # --- Header (light bar + title) ---
    banner = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.25), Inches(0.12), Inches(12.85), Inches(1.12)
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = RGBColor(0xFF, 0xF5, 0xF5)
    banner.line.fill.background()
    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.15), Inches(12.5), Inches(1.05))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Parallel Weighted A*"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xCC, 0x00, 0x00)
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "Jillian Chong (jc2886) · Julian Park (jp2534) · Lina Hao (lh657)"
    p2.font.size = Pt(14)
    p2.alignment = PP_ALIGN.CENTER
    p3 = tf.add_paragraph()
    p3.text = "Cornell University · CS 5220"
    p3.font.size = Pt(12)
    p3.font.italic = True
    p3.alignment = PP_ALIGN.CENTER

    intro = (
        "We explore a hybrid MPI/OpenMP parallel implementation of weighted A*. Weighted A* scales the "
        "heuristic in f = g + w·h to trade optimality for fewer expansions when h is informative. "
        "We vary MPI ranks, OpenMP threads per rank, heuristic weight w (vs. serial), and graph families: "
        "implicit grids with obstacles, geometric k-nearest-neighbor graphs, and CSR edge lists. "
        "Hypothesis: larger w narrows the frontier and can reduce cross-rank work—subject to graph structure "
        "and distributed ordering effects."
    )

    findings = (
        "• w > 1 often cuts expansions vs w = 1 on grid/geom (suboptimal paths possible).\n"
        "• MPI can speed large geom-like workloads; on easy grids MPI may trail serial (ordering + messaging).\n"
        "• OpenMP mainly helps high-degree neighbor work; grids/geom usually expand neighbors serially.\n"
        "• Many-thread ranks contend on locked frontier updates—1 thread/rank is often competitive."
    )

    methods = (
        "• Algorithm: distributed supersteps; each rank keeps local OPEN/gbest; relaxations routed by vertex "
        "partition (contiguous IDs).\n"
        "• MPI: batched exchange of relax messages between ranks (collective-style all-to-all pattern).\n"
        "• OpenMP: optional parallel neighbor staging when degree is high; parallel scan of inbound batches "
        "with serialized updates to shared frontier state.\n"
        "• Baselines: serial vs hybrid; weighted A* via flag -w; implicit graphs avoid storing 10^10–10^11 edges explicitly."
    )

    analyses = (
        "• Parallel benefit needs sufficient work per step; narrow frontiers favor serial A* ordering.\n"
        "• Distributed search ≠ textbook serial A*: expansion counts and even reported costs can diverge.\n"
        "• Higher w reduces greedy conservatism but only where h ≠ 0 (CSR with h ≡ 0: w unused for ordering)."
    )

    results = (
        "Plot timing vs scenario (serial vs MPI layouts: ranks × threads). Line charts & tables:\n"
        "results/w1_1e10_lines.(svg|jpg), results/w5_1e11_lines.(svg|jpg); raw TSV alongside.\n"
        "Patterns: strong MPI speedups on heavy geom at large n; many implicit-grid cases show serial competitive "
        "or faster program time_s; w = 5 reduces expansions sharply vs w = 1 on grids."
    )

    conclusions = (
        "Weighted w is an effective knob when h > 0. Hybrid MPI/OpenMP helps heavy searches but is not universally "
        "faster than serial—validate per graph family and (ranks, threads, w). Prefer profiling frontier updates "
        "and messaging before scaling ranks."
    )

    references = (
        "Hart, Nilsson, Raphael—A* (1968). Weighted / ε-admissible A* literature. MPI & OpenMP specifications "
        "(versions used). Course materials—CS 5220."
    )

    ack = "Thanks to CS 5220 staff and collaborators. Acknowledge compute facility / allocation if applicable."

    y0 = Inches(1.28)
    col_w = Inches(4.05)
    gap = Inches(0.15)
    left_col = Inches(0.35)
    mid_col = left_col + col_w + gap
    right_col = mid_col + col_w + gap

    # Left column (boxed)
    add_card(slide, left_col, y0, col_w, Inches(2.35), "Introduction", intro)
    add_card(slide, left_col, y0 + Inches(2.42), col_w, Inches(1.95), "Key Findings", findings)
    add_card(slide, left_col, y0 + Inches(4.45), col_w, Inches(2.68), "Methods", methods)

    # Middle column
    add_card(slide, mid_col, y0, col_w, Inches(2.15), "Analyses", analyses)
    add_card(
        slide,
        mid_col,
        y0 + Inches(2.25),
        col_w,
        Inches(4.85),
        "Results",
        results,
        title_size=13,
        body_size=9,
    )

    # Right column
    add_card(slide, right_col, y0, col_w, Inches(3.35), "Conclusions", conclusions)
    add_card(slide, right_col, y0 + Inches(3.45), col_w, Inches(1.55), "References", references, body_size=8)
    add_card(slide, right_col, y0 + Inches(5.08), col_w, Inches(2.05), "Acknowledgments", ack)

    out = "docs/CS5220_Parallel_Weighted_Astar_poster.pptx"
    prs.save(out)
    print("Wrote", out)


if __name__ == "__main__":
    main()
